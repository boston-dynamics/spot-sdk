# Copyright (c) 2023 Boston Dynamics, Inc.  All rights reserved.
#
# Downloading, reproducing, distributing or otherwise using the SDK Software
# is subject to the terms and conditions of the Boston Dynamics Software
# Development Kit License (20191101-BDSDK-SL).

import argparse
import os
import shutil
import sys
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation

from bosdyn.api import image_pb2
from bosdyn.api.autowalk import walks_pb2


def read_walk(walk_file):
    """ Function to read out the walk proto from a file into a walks_pb2
    - Args:
        - walk_file (Path): A .walk file generated from a recorded mission
    - Returns:
        - walk (walks_pb2): A walks_pb2 containing autowalk proto from the input file
    """
    walk = walks_pb2.Walk()

    with open(walk_file, 'rb') as autowalk_file:
        data = autowalk_file.read()
        walk.ParseFromString(data)

    return walk


def save_image(mission_image, image_location):
    """ Function to find format of image and save it to a file.
    - Args:
        - mission_image (image_pb2): The image to be saved
        - image_location (Path): File location for the saved image
    - Returns:
        - boolean: True if image is succesfully saved False otherwise.
    """
    if mission_image.format == image_pb2.Image.FORMAT_JPEG:
        image_location.write_bytes(mission_image.data)
        return True
    # Handling for thermal and depth images see https://dev.bostondynamics.com/docs/concepts/data_acquisition_thermal_raw.html for details
    elif (mission_image.format == image_pb2.Image.FORMAT_RAW):
        if mission_image.pixel_format not in (image_pb2.Image.PIXEL_FORMAT_DEPTH_U16,
                                              image_pb2.Image.PIXEL_FORMAT_GREYSCALE_U16):
            print('Warning, unsupported RAW format')
            return False
        # Thermal and depth data is raw 16 bit unsigned values
        image_array = np.frombuffer(mission_image.data, dtype='uint16')
        # Resize array to image size
        image_array = np.reshape(image_array, (mission_image.rows, mission_image.cols))
        # Plot and save the thermal data
        plt.imsave(image_location, image_array, cmap=plt.cm.inferno)
        return True
    elif mission_image.format == image_pb2.Image.FORMAT_RLE:
        print('Warning, RLE format not supported.')
        return False
    elif mission_image.format == image_pb2.Image.FORMAT_UNKNOWN:
        print('Warning image with unknown format')
        return False


def insert_next_slide(presentation, title, caption, image_location, mission_image):
    """ Function to insert the next slide into a presentation with an image and caption
    - Args:
        - presentation (Presentation): Variable containing the presentation to be added to
        - title (string): Title text for the slide
        - caption (string): Caption text for the slide
        - image_location (Path): File location for the image to be inserted
        - mission_image (image_pb2): image_pb2 file of the image to be inserted
    """
    # This is a built-in format for a slide with an image and caption
    slide = presentation.slides.add_slide(presentation.slide_layouts[8])
    slide.placeholders[0].text = title
    # In this format the caption text always in position 2
    slide.placeholders[2].text = caption
    # In this format the image placeholder is always in position 1
    imagePlaceholder = slide.placeholders[1]
    with open(image_location, 'rb') as file:
        picture = imagePlaceholder.insert_picture(file)
    # Height and width values disappear after modifying one, we save the original ones so we don't distort the image
    old_width = picture.width
    old_height = picture.height
    old_top = picture.top
    old_left = picture.left
    picture.crop_top = 0
    picture.crop_left = 0
    picture.crop_right = 0
    picture.crop_bottom = 0
    scale_factor = (mission_image.cols / mission_image.rows) / (old_width / old_height)
    # Always scale by expanding width since that's where we have more room on the slide.
    picture.width = round(old_width * scale_factor)
    picture.height = old_height
    picture.left = old_left - 500000  # Ballpark number that makes most images centered
    picture.top = old_top


def prepare_directories(options):
    """ Function to prepare the output directories for the presentation and images. Deletes directories if they
    already exist, creates them otherwise.
    - Args:
        - options (Namespace): Variable containing the parsed arguments passed in from running this program.
    - Returns:
        - slides_dir (Path): Directory to save the output presentation slides
        - images_dir (Path): Directory to save the output images
    """
    images_dir = options.output_directory / 'images'
    if images_dir.exists():
        shutil.rmtree(images_dir)

    slides_dir = options.output_directory / 'slides'
    if slides_dir.exists():
        shutil.rmtree(slides_dir)

    images_dir.mkdir(parents=True)
    slides_dir.mkdir(parents=True)

    return slides_dir, images_dir


def process_walk(slides_dir, images_dir, walk):
    """ Function to process walk proto into slide and image outputs, saved as described in the README
    - Args:
        - slides_dir (Path): Directory to save the output presentation slides
        - images_dir (Path): Directory to save the output images
        - walk (walks_pb2): A walks_pb2 to be processed
    """
    presentation = Presentation()
    num_element_processing = 0
    for element in walk.elements:  # Iterate through all elements in walk
        num_element_processing += 1
        print(f'Processing {num_element_processing} of {len(walk.elements)}')
        # Only care about elements with DAQ actions with images
        if len(element.action.data_acquisition.record_time_images) >= 1:
            num_record_time_image = 0
            for record_time_image in element.action.data_acquisition.record_time_images:  # Iterate through all images within an action
                image_location = images_dir / f'{element.name}_{num_record_time_image}.jpeg'

                if not save_image(record_time_image.shot.image, image_location):
                    continue

                caption = record_time_image.source.name
                # Special handling for Soundsurface actions to show the frequency range
                if caption == 'Soundsurface':
                    spec = record_time_image.source.custom_params.specs[
                        "Frequency From"].spec.int_spec

                    caption += (f'\n Default min freq: {spec.default_value.value} Hz \n'
                                f'min_value: {spec.min_value.value} Hz\n'
                                f'max_value: {spec.max_value.value} Hz')

                # If we have any network compute captures for this action we add the config(s) to our caption
                for networkcomp in element.action.data_acquisition.acquire_data_request.acquisition_requests.network_compute_captures:
                    caption += f'\n{str(networkcomp.server_config)}'

                insert_next_slide(presentation, element.name, caption, image_location,
                                  record_time_image.shot.image)
                num_record_time_image += 1

    presentation.save(slides_dir / 'Output.pptx')


def main():
    parser = argparse.ArgumentParser()

    parser.add_argument('--input_file', type=Path,
                        help='Path to the .walk file to extract images from', required=True)

    parser.add_argument('--output_directory', default='', type=Path,
                        help='Path to the directory to output to, default is current directory')

    options = parser.parse_args()
    slides_dir, images_dir = prepare_directories(options)
    walk = read_walk(options.input_file)
    process_walk(slides_dir, images_dir, walk)


if __name__ == '__main__':
    if not main():
        sys.exit(1)
